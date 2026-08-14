"""
Page-aware loader for the "page" splitter strategy.

MarkItDown returns one flat text string, and by the time a splitter sees
that string the page boundaries are gone. This loader keeps them: PDFs are
extracted page by page with pypdf, PPTX decks slide by slide with
python-pptx. Text-form files are returned raw so PageSplitter can split
them on the page marker or on headings, and every other format falls back
to MarkItDown as a single flat text.
"""

import logging
from pathlib import Path
from typing import Any, Dict, List, Union

from .markitdown_loader import MarkItDownLoader

logger = logging.getLogger(__name__)


class PageLoader:
    """
    Load documents with page boundaries preserved.

    Returns the same dict as MarkItDownLoader plus a ``pages`` key: a list
    of ``{"page_number", "page_label", "text"}`` dicts for formats with real
    pages (PDF, PPTX), or None when finding page boundaries is the
    splitter's job (markers in text files, markdown headings) or pagination
    does not apply.

    Example:
        >>> loader = PageLoader()
        >>> result = loader.load("report.pdf")
        >>> result["pages"][0]["page_number"]
        1
    """

    #: Read directly from disk; PageSplitter finds the page boundaries.
    RAW_TEXT_TYPES = {"txt", "md", "markdown", "text"}

    def __init__(self):
        """Initialize the loader and its MarkItDown fallback."""
        self.fallback = MarkItDownLoader()

    def load(self, file_path: Union[str, Path]) -> dict:
        """
        Load one document, preserving page boundaries where the format has them.

        Args:
            file_path: Path to the document file

        Returns:
            Dictionary containing:
                - text_content: Full extracted text
                - file_name: Original filename
                - file_type: File extension
                - success: Whether extraction succeeded
                - error: Error message when it did not
                - pages: List of page dicts, or None (see class docstring)

        Raises:
            FileNotFoundError: If file doesn't exist
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_type = file_path.suffix.lower().lstrip(".")

        try:
            if file_type == "pdf":
                pages = self._load_pdf(file_path)
            elif file_type == "pptx":
                pages = self._load_pptx(file_path)
            elif file_type in self.RAW_TEXT_TYPES:
                text = file_path.read_text(encoding="utf-8", errors="replace")
                return self._result(file_path, file_type, text, pages=None)
            else:
                # No page concept (DOCX has none until rendered); flat text.
                result = self.fallback.load(file_path)
                result["pages"] = None
                return result
        except Exception as e:
            logger.error(f"Failed to load pages from {file_path}: {e}")
            return {
                "text_content": "",
                "file_name": file_path.name,
                "file_type": file_type,
                "success": False,
                "error": str(e),
                "pages": None,
            }

        text_content = "\n\n".join(page["text"] for page in pages)
        return self._result(file_path, file_type, text_content, pages=pages)

    @staticmethod
    def _result(
        file_path: Path, file_type: str, text: str, pages
    ) -> Dict[str, Any]:
        return {
            "text_content": text,
            "file_name": file_path.name,
            "file_type": file_type,
            "success": True,
            "error": None,
            "pages": pages,
        }

    @staticmethod
    def _load_pdf(file_path: Path) -> List[Dict[str, Any]]:
        """
        One entry per physical PDF page, in order.

        The page label carries the document's printed numbering ("iv",
        "A-1") only when the PDF defines one that differs from the ordinal,
        so a plain 1, 2, 3 sequence adds no redundant field.
        """
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))

        try:
            labels = reader.page_labels
        except Exception:
            labels = None

        pages = []
        for index, page in enumerate(reader.pages):
            label = None
            if labels is not None and index < len(labels):
                if str(labels[index]) != str(index + 1):
                    label = str(labels[index])
            pages.append(
                {
                    "page_number": index + 1,
                    "page_label": label,
                    "text": page.extract_text() or "",
                }
            )
        return pages

    @staticmethod
    def _load_pptx(file_path: Path) -> List[Dict[str, Any]]:
        """
        One entry per slide; the slide title becomes the page label.

        Text is gathered from every text frame on the slide, and table cells
        are flattened row by row so tabular decks stay searchable.
        """
        from pptx import Presentation

        presentation = Presentation(str(file_path))

        pages = []
        for index, slide in enumerate(presentation.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
                elif getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.append(
                            " | ".join(cell.text.strip() for cell in row.cells)
                        )

            title = None
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape is not None and title_shape.text.strip():
                title = title_shape.text.strip()

            pages.append(
                {
                    "page_number": index + 1,
                    "page_label": title,
                    "text": "\n".join(parts),
                }
            )
        return pages
